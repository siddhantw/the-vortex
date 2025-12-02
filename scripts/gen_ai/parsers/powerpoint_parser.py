from .base_parser import BaseParser
from pptx import Presentation

class PowerPointParser(BaseParser):
    def parse(self, file_path: str) -> dict:
        text = ""
        try:
            # Load the PowerPoint presentation
            prs = Presentation(file_path)

            # Extract text from each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"=== Slide {slide_num} ===\n"

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text += " | ".join(row_text) + "\n"

                text += "\n"  # Add spacing between slides

        except Exception as e:
            text = f"[PowerPoint parsing error: {e}]"

        return {"raw_text": text}
